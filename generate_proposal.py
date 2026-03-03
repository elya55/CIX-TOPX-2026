import sys
sys.path.insert(0, '/Users/ccelya/Downloads/AI-SVW-C/lib')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─────────────────────────────────────────────
# 品牌色彩系统
# ─────────────────────────────────────────────
BG_DARK      = RGBColor(0x0A, 0x0A, 0x14)   # 深夜蓝黑 - 主背景
BG_MID       = RGBColor(0x10, 0x10, 0x22)   # 次深背景
ACCENT_BLUE  = RGBColor(0x00, 0x8B, 0xFF)     # 科技蓝主色
ACCENT_LIGHT = RGBColor(0x4D, 0xB8, 0xFF)   # 浅蓝高光
ACCENT_GRAD  = RGBColor(0x00, 0xD4, 0xAA)   # 青绿渐变色
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_70     = RGBColor(0xB3, 0xB3, 0xCC)   # 70%白
WHITE_40     = RGBColor(0x66, 0x66, 0x88)   # 40%白
GOLD         = RGBColor(0xFF, 0xC8, 0x00)   # 金色强调
RED_LIGHT    = RGBColor(0xFF, 0x4D, 0x6D)   # 红色警示
TEAL         = RGBColor(0x00, 0xC9, 0xA7)   # 青色

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
blank_layout = prs.slide_layouts[6]

# ─────────────────────────────────────────────
# 辅助函数
# ─────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=BG_DARK):
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def rect(slide, x, y, w, h, fill_color, alpha=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        font_size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def txt_multiline(slide, lines, x, y, w, h,
                  font_size=Pt(13), bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, line_spacing=None):
    """lines: list of (text, bold, color, size) or just str"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, bold, color, font_size)
        t, b, c, s = item[0], item[1] if len(item)>1 else bold, item[2] if len(item)>2 else color, item[3] if len(item)>3 else font_size
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = t
        run.font.size = s
        run.font.bold = b
        run.font.color.rgb = c
    return txBox

def divider_line(slide, x, y, w, color=ACCENT_BLUE, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(1, x, y, w, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line

def tag_badge(slide, text, x, y, bg=ACCENT_BLUE, fg=WHITE, size=Pt(11)):
    w = Inches(1.5)
    h = Inches(0.32)
    r = rect(slide, x, y, w, h, bg)
    txt(slide, text, x, y + Inches(0.04), w, h, font_size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return r

def chapter_divider(slide, number, title, subtitle=""):
    fill_bg(slide, BG_DARK)
    # 左侧大色块
    rect(slide, 0, 0, Inches(4.5), H, RGBColor(0x00, 0x55, 0xCC))
    # 章节号
    txt(slide, f"0{number}", Inches(0.5), Inches(1.8), Inches(3.5), Inches(2),
        font_size=Pt(120), bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    # 标题
    txt(slide, title, Inches(5), Inches(2.8), Inches(7.8), Inches(1.2),
        font_size=Pt(36), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(5), Inches(4.1), Inches(7.8), Inches(0.8),
            font_size=Pt(18), color=WHITE_70)
    divider_line(slide, Inches(5), Inches(2.6), Inches(7), ACCENT_BLUE, Pt(3))

def section_header(slide, title, subtitle=None, tag=None):
    """每页左上角小标题区域"""
    if tag:
        tag_badge(slide, tag, Inches(0.4), Inches(0.3))
        txt(slide, title, Inches(0.4), Inches(0.7), Inches(8), Inches(0.55),
            font_size=Pt(22), bold=True, color=WHITE)
    else:
        txt(slide, title, Inches(0.4), Inches(0.3), Inches(8), Inches(0.55),
            font_size=Pt(22), bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(1.0), Inches(10), Inches(0.4),
            font_size=Pt(13), color=WHITE_70)
    divider_line(slide, Inches(0.4), Inches(1.4), Inches(12.5), ACCENT_BLUE)

def card(slide, x, y, w, h, title, body_lines, title_color=ACCENT_LIGHT, icon=""):
    r = rect(slide, x, y, w, h, RGBColor(0x16, 0x16, 0x2E), line_color=RGBColor(0x1E,0x1E,0x44), line_width=Pt(1))
    # 顶部色条
    rect(slide, x, y, w, Inches(0.06), ACCENT_BLUE)
    ty = y + Inches(0.15)
    if icon:
        txt(slide, icon, x + Inches(0.15), ty, Inches(0.4), Inches(0.4), font_size=Pt(18), color=ACCENT_BLUE)
        txt(slide, title, x + Inches(0.55), ty, w - Inches(0.7), Inches(0.4),
            font_size=Pt(14), bold=True, color=title_color)
    else:
        txt(slide, title, x + Inches(0.2), ty, w - Inches(0.4), Inches(0.4),
            font_size=Pt(14), bold=True, color=title_color)
    by = ty + Inches(0.45)
    for line in body_lines:
        if isinstance(line, tuple):
            t, c, s = line[0], line[1] if len(line)>1 else WHITE_70, line[2] if len(line)>2 else Pt(12)
        else:
            t, c, s = line, WHITE_70, Pt(12)
        txt(slide, t, x + Inches(0.2), by, w - Inches(0.4), Inches(0.35),
            font_size=s, color=c)
        by += Inches(0.32)

def bullet_list(slide, items, x, y, w, color=WHITE_70, size=Pt(13), bullet="▸ "):
    for item in items:
        txt(slide, bullet + item, x, y, w, Inches(0.35), font_size=size, color=color)
        y += Inches(0.38)
    return y

# ═══════════════════════════════════════════════════════
# SLIDE 1: 封面
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)

# 左侧大渐变蓝条
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)

# 顶部细线
rect(s, Inches(0.3), Inches(0.4), Inches(12.7), Inches(0.02), ACCENT_BLUE)

# 大标题
txt(s, "2025 上汽大众", Inches(0.5), Inches(1.2), Inches(12), Inches(0.9),
    font_size=Pt(46), bold=True, color=WHITE)
txt(s, "C端产品 & C车 UIUX 设计提案", Inches(0.5), Inches(2.1), Inches(12), Inches(0.9),
    font_size=Pt(38), bold=True, color=ACCENT_BLUE)

# 副标题
txt(s, "SVW C-Touch & CMP21D Vehicle App | Full-Stack UX/UI Design Proposal",
    Inches(0.5), Inches(3.1), Inches(12), Inches(0.5),
    font_size=Pt(16), color=WHITE_70, italic=True)

divider_line(s, Inches(0.5), Inches(3.7), Inches(8), ACCENT_BLUE, Pt(2))

# 标签组
for i, (label, val) in enumerate([
    ("项目方", "上汽大众汽车有限公司 CIX"),
    ("提案方", "TOPX Design"),
    ("日期", "2025.11.14"),
]):
    x = Inches(0.5 + i * 4.0)
    txt(s, label, x, Inches(4.0), Inches(1.5), Inches(0.3),
        font_size=Pt(11), color=WHITE_40)
    txt(s, val, x, Inches(4.3), Inches(3.8), Inches(0.4),
        font_size=Pt(14), bold=True, color=WHITE)

# 底部装饰
rect(s, 0, H - Inches(0.5), W, Inches(0.5), RGBColor(0x00, 0x33, 0x88))
txt(s, "CONFIDENTIAL  |  TOPX Design 2025  |  For SVW CIX Only",
    Inches(0.5), H - Inches(0.42), Inches(12), Inches(0.35),
    font_size=Pt(10), color=WHITE_40, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 2: 目录
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)

txt(s, "目 录", Inches(0.5), Inches(0.4), Inches(4), Inches(0.7),
    font_size=Pt(30), bold=True, color=WHITE)
txt(s, "Contents", Inches(0.5), Inches(1.1), Inches(4), Inches(0.4),
    font_size=Pt(14), color=WHITE_40, italic=True)
divider_line(s, Inches(0.5), Inches(1.55), Inches(12.3), ACCENT_BLUE)

chapters = [
    ("01", "关于我们 About Us", "公司介绍 · 核心团队 · 方法论 · 合作案例", ACCENT_BLUE),
    ("02", "背景分析 Background", "战略解读 · 行业趋势 · 竞品分析 · 现状诊断", ACCENT_LIGHT),
    ("03A", "C端设计方案", "官网 · 小程序 · App 爱车 · 视觉风格", ACCENT_GRAD),
    ("03B", "C车设计方案", "CMP21D 车控 App · 关键功能 · 双品牌融合", TEAL),
    ("04", "项目计划 Plan", "团队架构 · 实施规划 · 时间表 · 交付物", GOLD),
]

for i, (num, title, sub, color) in enumerate(chapters):
    y = Inches(1.8 + i * 1.0)
    rect(s, Inches(0.5), y, Inches(0.7), Inches(0.65), color)
    txt(s, num, Inches(0.5), y + Inches(0.1), Inches(0.7), Inches(0.5),
        font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.4), y + Inches(0.05), Inches(5), Inches(0.35),
        font_size=Pt(18), bold=True, color=WHITE)
    txt(s, sub, Inches(1.4), y + Inches(0.38), Inches(9), Inches(0.3),
        font_size=Pt(12), color=WHITE_40)
    divider_line(s, Inches(1.4), y + Inches(0.7), Inches(11.4), RGBColor(0x22,0x22,0x44), Pt(0.5))

# ═══════════════════════════════════════════════════════
# SLIDE 3: 章节分隔 - 01 About Us
# ═══════════════════════════════════════════════════════
s = add_slide()
chapter_divider(s, 1, "关于我们", "About Us  ·  Who We Are")

# ═══════════════════════════════════════════════════════
# SLIDE 4: 公司介绍
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "TOPX Design 公司介绍", "设计创造价值 · Design Creates Value", tag="01 About Us")

txt(s, '"设计创造价值"', Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.7),
    font_size=Pt(26), bold=True, color=ACCENT_BLUE)
txt(s, "Design Creates Value",
    Inches(0.5), Inches(2.35), Inches(5.5), Inches(0.4),
    font_size=Pt(14), color=WHITE_40, italic=True)

desc = ("TOPX 成立于 2017 年 5 月，由百度、腾讯等互联网背景的成员组成，"
        "获真格基金投资，是国内领先的创新型用户体验及用户增长咨询机构。"
        "深度服务大众、奥迪、保时捷、林肯等顶级汽车品牌数字化设计项目。")
txt(s, desc, Inches(0.5), Inches(2.85), Inches(5.8), Inches(1.4),
    font_size=Pt(13), color=WHITE_70)

# 右侧数据墙
stats = [
    ("U9", "最高级别体验设计师", ACCENT_BLUE),
    ("10年+", "用户体验行业经验", ACCENT_LIGHT),
    ("20+", "优质客户 多次合作", ACCENT_GRAD),
    ("100+", "创新设计方案", TEAL),
    ("4", "分公司：北京/上海/杭州/合肥", GOLD),
    ("90+", "线上线下行业分享", WHITE_70),
]
for i, (num, label, color) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = Inches(6.5 + col * 2.2)
    y = Inches(1.7 + row * 1.4)
    rect(s, x, y, Inches(2.0), Inches(1.2), RGBColor(0x12,0x12,0x28),
         line_color=RGBColor(0x22,0x22,0x44), line_width=Pt(1))
    rect(s, x, y, Inches(2.0), Inches(0.05), color)
    txt(s, num, x, y + Inches(0.1), Inches(2.0), Inches(0.55),
        font_size=Pt(28), bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, x, y + Inches(0.62), Inches(2.0), Inches(0.5),
        font_size=Pt(10), color=WHITE_70, align=PP_ALIGN.CENTER, wrap=True)

# 大众经验强调
rect(s, Inches(0.5), Inches(4.8), Inches(5.8), Inches(0.6), RGBColor(0x00,0x33,0x88))
txt(s, "🏆  大众·奥迪·斯柯达·保时捷·林肯 数字化设计项目核心负责人",
    Inches(0.7), Inches(4.85), Inches(5.4), Inches(0.5),
    font_size=Pt(12), bold=True, color=WHITE)

# ═══════════════════════════════════════════════════════
# SLIDE 5: 核心团队
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "核心项目团队", "专业配置 · 对应 RFP 全岗位要求", tag="01 About Us")

roles = [
    ("Elya", "设计总监\n项目经理", "14年体验设计 / TOPX CEO\n前百度云UED负责人\n大众·保时捷·奥迪项目负责人", ACCENT_BLUE),
    ("Maya", "创意总监\nDesign Leader", "15年产品设计 / TOPX CDO\n前三星·百度MUX资深交互\n保时捷·林肯·现代项目负责人", ACCENT_LIGHT),
    ("Arthur", "产品总监\nStrategy Leader", "15年顶层产品设计 / TOPX CMO\n前独角兽用户增长负责人\n极氪·吉利·奥迪全链路专家", ACCENT_GRAD),
    ("James", "UI设计总监", "2021红点奖获得者\n高级UI设计 10年+经验\n车企App视觉专家", TEAL),
    ("Kay", "高级设计经理", "11年从业经验\n高级交互/视觉设计\n车联网App专家", GOLD),
    ("Mira", "高级用户研究员", "用户调研专家\n可用性测试 / 焦点小组\n用户画像建模专家", WHITE_70),
]

for i, (name, role, exp, color) in enumerate(roles):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.65 + row * 2.5)
    rect(s, x, y, Inches(4.0), Inches(2.2), RGBColor(0x12,0x12,0x2A),
         line_color=RGBColor(0x22,0x22,0x44), line_width=Pt(1))
    rect(s, x, y, Inches(0.5), Inches(2.2), color)
    txt(s, name, x + Inches(0.65), y + Inches(0.15), Inches(3.2), Inches(0.5),
        font_size=Pt(18), bold=True, color=WHITE)
    for j, line in enumerate(role.split("\n")):
        txt(s, line, x + Inches(0.65), y + Inches(0.62 + j*0.32), Inches(3.2), Inches(0.32),
            font_size=Pt(12), bold=True, color=color)
    for j, line in enumerate(exp.split("\n")):
        txt(s, "· " + line, x + Inches(0.65), y + Inches(1.2 + j*0.3), Inches(3.2), Inches(0.3),
            font_size=Pt(10.5), color=WHITE_40)

# ═══════════════════════════════════════════════════════
# SLIDE 6: 方法论
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "TOP3 独家方法论体系", "5维体验诊断 · STAR规划 · 增长模型", tag="01 About Us")

# 5维体验诊断
txt(s, "5维体验诊断", Inches(0.5), Inches(1.65), Inches(4), Inches(0.45),
    font_size=Pt(16), bold=True, color=ACCENT_BLUE)
layers = [
    ("增长层", GOLD), ("诊断层", ACCENT_BLUE), ("战略层", ACCENT_LIGHT),
    ("体验层", ACCENT_GRAD), ("结构层", TEAL),
]
for i, (name, color) in enumerate(layers):
    w_val = Inches(3.5 - i * 0.3)
    x_val = Inches(0.5 + i * 0.15)
    rect(s, x_val, Inches(2.15 + i * 0.65), w_val, Inches(0.55), color)
    txt(s, name, x_val + Inches(0.15), Inches(2.2 + i * 0.65), w_val, Inches(0.45),
        font_size=Pt(13), bold=True, color=WHITE)

# STAR规划
txt(s, "STAR 实施规划", Inches(4.8), Inches(1.65), Inches(4), Inches(0.45),
    font_size=Pt(16), bold=True, color=ACCENT_BLUE)
stars = [
    ("S", "Situation 现状诊断", "专家走查 · 竞品分析 · 用户访谈", ACCENT_BLUE),
    ("T", "Task 需求分析", "KANO模型 · SWOT · 行业趋势", ACCENT_LIGHT),
    ("A", "Action 方案设计", "信息架构 · 交互设计 · 视觉设计", ACCENT_GRAD),
    ("R", "Result 结果验证", "可用性测试 · A/B测试 · 数据迭代", GOLD),
]
for i, (letter, title, desc, color) in enumerate(stars):
    y = Inches(2.15 + i * 1.12)
    rect(s, Inches(4.8), y, Inches(0.55), Inches(0.95), color)
    txt(s, letter, Inches(4.8), y + Inches(0.2), Inches(0.55), Inches(0.55),
        font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(5.5), y + Inches(0.1), Inches(4.2), Inches(0.35),
        font_size=Pt(13), bold=True, color=color)
    txt(s, desc, Inches(5.5), y + Inches(0.47), Inches(4.2), Inches(0.35),
        font_size=Pt(11), color=WHITE_40)

# AARRR增长模型
txt(s, "5维增长模型 AARRR", Inches(10.2), Inches(1.65), Inches(2.8), Inches(0.45),
    font_size=Pt(14), bold=True, color=ACCENT_BLUE)
aarrr = [
    ("A", "Acquisition 拉新", ACCENT_BLUE),
    ("A", "Activation 促活", ACCENT_LIGHT),
    ("R", "Retention 留存", ACCENT_GRAD),
    ("R", "Revenue 变现", TEAL),
    ("R", "Referral 传播", GOLD),
]
for i, (letter, label, color) in enumerate(aarrr):
    y = Inches(2.2 + i * 0.88)
    rect(s, Inches(10.2), y, Inches(0.4), Inches(0.7), color)
    txt(s, letter, Inches(10.2), y + Inches(0.15), Inches(0.4), Inches(0.4),
        font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, Inches(10.75), y + Inches(0.18), Inches(2.3), Inches(0.38),
        font_size=Pt(12), color=WHITE_70)

# ═══════════════════════════════════════════════════════
# SLIDE 7: 大众合作历史
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "与上汽大众的深厚渊源", "5年+ 持续合作 · 深度理解品牌DNA", tag="01 About Us")

txt(s, "与大众的缘分…源自 5 年前",
    Inches(0.5), Inches(1.7), Inches(12), Inches(0.5),
    font_size=Pt(20), bold=True, color=WHITE)

projects = [
    ("2019", "VOLKSWAGEN HUX", "上汽大众电商平台体验重塑", ACCENT_BLUE, "官网体验 · 电商转化 · 用户研究"),
    ("2020", "ID. 系列在线展厅", "上汽ID.系列新车线上体验设计", ACCENT_LIGHT, "沉浸式展厅 · 车型浏览 · 配车流程"),
    ("2021-22", "大众 App 运营", "智选充电/ID电力挑战赛/爱奇艺/云听等", ACCENT_GRAD, "运营活动设计 · 会员体系 · 促活留存"),
    ("2023", "C端产品 UIUX", "官网+小程序+App 五维体验重塑", TEAL, "官网重构 · 情感化设计 · Design System"),
    ("2024-25", "AB车 UIUX", "A/B车型车控App 全场景情感化设计", GOLD, "车控设计 · Avatar · 大模型融合"),
    ("2025 →", "C端 & C车 (新)", "ID.ERA + CMP21D 双项目全链路设计", RED_LIGHT, "本次提案项目"),
]

for i, (year, title, desc, color, tag_txt) in enumerate(projects):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(2.4 + row * 2.1)
    rect(s, x, y, Inches(4.0), Inches(1.85), RGBColor(0x10,0x10,0x26),
         line_color=color, line_width=Pt(1))
    rect(s, x, y, Inches(0.9), Inches(0.38), color)
    txt(s, year, x, y + Inches(0.06), Inches(0.9), Inches(0.3),
        font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x + Inches(1.0), y + Inches(0.08), Inches(2.8), Inches(0.3),
        font_size=Pt(13), bold=True, color=color)
    txt(s, desc, x + Inches(0.15), y + Inches(0.55), Inches(3.7), Inches(0.45),
        font_size=Pt(11.5), color=WHITE_70, wrap=True)
    txt(s, tag_txt, x + Inches(0.15), y + Inches(1.1), Inches(3.7), Inches(0.55),
        font_size=Pt(10.5), color=WHITE_40, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 8: 章节分隔 - 02 背景分析
# ═══════════════════════════════════════════════════════
s = add_slide()
chapter_divider(s, 2, "背景分析", "Background Analysis  ·  Strategy & Insights")

# ═══════════════════════════════════════════════════════
# SLIDE 9: 战略解读
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "战略背景解读", "双项目协同 · 构建大众中国数字体验生态", tag="02 背景分析")

# 两个项目背景卡片
for col, (proj, color, bg_c, lines) in enumerate([
    ("C端项目", ACCENT_BLUE, RGBColor(0x00,0x22,0x55), [
        ("ID.ERA 增程电动 SUV", WHITE, Pt(15), True),
        ("2026 年首款增程电动 SUV", WHITE_70, Pt(12), False),
        ('"In China, for China" 战略核心产品', ACCENT_LIGHT, Pt(12), False),
        ("覆盖：官网 · 小程序 · App 爱车", WHITE_70, Pt(12), False),
        ("目标：潜在/存量/增量用户全链路体验", WHITE_70, Pt(12), False),
    ]),
    ("C车项目", TEAL, RGBColor(0x00,0x33,0x2A), [
        ("CMP21D 大众 × 小鹏联合车型", WHITE, Pt(15), True),
        ("2026 年首款联合开发新车型", WHITE_70, Pt(12), False),
        ("融合大众可信赖 + 小鹏智能化基因", ACCENT_GRAD, Pt(12), False),
        ("覆盖：车主专属移动互联 App", WHITE_70, Pt(12), False),
        ("目标：极致智能车控体验", WHITE_70, Pt(12), False),
    ]),
]):
    x = Inches(0.4 + col * 6.4)
    rect(s, x, Inches(1.6), Inches(6.0), Inches(3.4), bg_c,
         line_color=color, line_width=Pt(1.5))
    rect(s, x, Inches(1.6), Inches(6.0), Inches(0.5), color)
    txt(s, proj, x + Inches(0.2), Inches(1.65), Inches(5.5), Inches(0.4),
        font_size=Pt(18), bold=True, color=WHITE)
    for j, (t, c, sz, b) in enumerate(lines):
        txt(s, t, x + Inches(0.25), Inches(2.25 + j * 0.45), Inches(5.5), Inches(0.4),
            font_size=sz, bold=b, color=c)

# 联动说明
rect(s, Inches(0.4), Inches(5.2), Inches(12.5), Inches(0.9), RGBColor(0x0A,0x1A,0x3A))
txt(s, "⚡  战略协同：", Inches(0.6), Inches(5.3), Inches(2), Inches(0.5),
    font_size=Pt(14), bold=True, color=ACCENT_BLUE)
txt(s, "C端负责「拉新获客」→ 打造品牌认知与购车转化 ｜ C车负责「用车留存」→ 构建车主粘性与口碑传播  ｜  共同构建大众车主全生命周期数字体验生态",
    Inches(2.6), Inches(5.32), Inches(10.2), Inches(0.5),
    font_size=Pt(12.5), color=WHITE_70)

# ═══════════════════════════════════════════════════════
# SLIDE 10: 行业趋势
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "行业趋势洞察", "左手：车企亮点  ·  右手：未来趋势  ·  融合探索", tag="02 背景分析")

txt(s, "左手：车企数字化亮点", Inches(0.5), Inches(1.65), Inches(5.5), Inches(0.4),
    font_size=Pt(15), bold=True, color=ACCENT_BLUE)
txt(s, "右手：行业未来趋势", Inches(7.1), Inches(1.65), Inches(5.5), Inches(0.4),
    font_size=Pt(15), bold=True, color=ACCENT_GRAD)

left_trends = [
    ("🔧 车控工具极致化", "远程控制 · 精准状态反馈 · 一键场景联动", ACCENT_BLUE),
    ("🤝 车控社交化", "车友圈 · 车主茶话会 · 用户共创生态", ACCENT_LIGHT),
    ("🧠 AI 大模型融合", "Avatar 智能体 · 语音克隆 · 情感识别", ACCENT_GRAD),
    ("🌐 C端沉浸体验", "一页式滚动 · 视差动效 · 场景分区构图", TEAL),
]
right_trends = [
    ("🚀 主动服务体验", "车随心动 · 行程预测 · 场景自动推荐", ACCENT_GRAD),
    ("❤️ 情感化设计升级", "香氛联动 · 情绪识别 · 个性化定制", GOLD),
    ("📱 多端生态整合", "手机/手表/车机 三端无缝协同体验", ACCENT_LIGHT),
    ("🔒 数据安全合规", "个人信息保护 · 车控安全 · 隐私设计", WHITE_70),
]

for i, (icon_title, desc, color) in enumerate(left_trends):
    y = Inches(2.15 + i * 1.1)
    rect(s, Inches(0.5), y, Inches(5.9), Inches(0.92), RGBColor(0x10,0x10,0x28),
         line_color=color, line_width=Pt(0.8))
    rect(s, Inches(0.5), y, Inches(0.06), Inches(0.92), color)
    txt(s, icon_title, Inches(0.7), y + Inches(0.08), Inches(5.5), Inches(0.35),
        font_size=Pt(14), bold=True, color=color)
    txt(s, desc, Inches(0.7), y + Inches(0.48), Inches(5.5), Inches(0.35),
        font_size=Pt(11.5), color=WHITE_40)

for i, (icon_title, desc, color) in enumerate(right_trends):
    y = Inches(2.15 + i * 1.1)
    rect(s, Inches(7.1), y, Inches(5.9), Inches(0.92), RGBColor(0x10,0x10,0x28),
         line_color=color, line_width=Pt(0.8))
    rect(s, Inches(7.1), y, Inches(0.06), Inches(0.92), color)
    txt(s, icon_title, Inches(7.3), y + Inches(0.08), Inches(5.5), Inches(0.35),
        font_size=Pt(14), bold=True, color=color)
    txt(s, desc, Inches(7.3), y + Inches(0.48), Inches(5.5), Inches(0.35),
        font_size=Pt(11.5), color=WHITE_40)

# 中间箭头
txt(s, "+", Inches(6.3), Inches(3.5), Inches(0.7), Inches(0.7),
    font_size=Pt(36), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
txt(s, "融合", Inches(6.2), Inches(4.2), Inches(0.9), Inches(0.4),
    font_size=Pt(12), color=WHITE_40, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 11: 竞品分析
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "竞品分析", "C端官网 & 车控App 双维度对标", tag="02 背景分析")

txt(s, "C端官网竞品", Inches(0.5), Inches(1.65), Inches(5.9), Inches(0.4),
    font_size=Pt(15), bold=True, color=ACCENT_BLUE)
txt(s, "车控 App 竞品", Inches(6.8), Inches(1.65), Inches(5.9), Inches(0.4),
    font_size=Pt(15), bold=True, color=TEAL)
divider_line(s, Inches(6.5), Inches(1.6), Inches(0.02), WHITE_40, Pt(1))

cend_comps = [
    ("Apple 官网", "✅ 沉浸式滚动·场景分区构图\n✅ 主体置中 RWD 响应式布局\n⚠️ 产品导向过强，情感触点少", RGBColor(0x60,0x60,0x60)),
    ("特斯拉官网", "✅ 简洁一页式·极简风格\n✅ 强 CTA 转化设计\n⚠️ 品牌温度感不足", RGBColor(0xCC,0x00,0x00)),
    ("蔚来官网", "✅ 情感化强·社区氛围浓厚\n✅ 用户旅程设计完整\n⚠️ 加载较重·转化效率待优化", RGBColor(0x00,0x99,0xFF)),
]
for i, (name, pros, color) in enumerate(cend_comps):
    y = Inches(2.1 + i * 1.6)
    rect(s, Inches(0.5), y, Inches(5.9), Inches(1.4), RGBColor(0x10,0x10,0x28),
         line_color=color, line_width=Pt(1))
    rect(s, Inches(0.5), y, Inches(1.5), Inches(0.4), color)
    txt(s, name, Inches(0.5), y + Inches(0.05), Inches(1.5), Inches(0.32),
        font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, line in enumerate(pros.split("\n")):
        txt(s, line, Inches(0.7), y + Inches(0.55 + j * 0.3), Inches(5.5), Inches(0.3),
            font_size=Pt(11), color=WHITE_70)

ccar_comps = [
    ("特斯拉 App", "✅ 远程控制强大·简洁高效\n✅ 工具属性极致化\n⚠️ 情感化设计缺失·体验单一", RGBColor(0xCC,0x00,0x00)),
    ("蔚来 App", "✅ 工具+转化+智能三维平衡\n✅ 车身状态实时同步\n⚠️ 功能整合便捷性待提升", RGBColor(0x00,0x99,0xFF)),
    ("小鹏 App", "✅ 智驾功能展示突出\n✅ AI 语音交互创新\n⚠️ 品牌视觉风格较弱", RGBColor(0x00,0xCC,0x88)),
]
for i, (name, pros, color) in enumerate(ccar_comps):
    y = Inches(2.1 + i * 1.6)
    rect(s, Inches(6.8), y, Inches(6.1), Inches(1.4), RGBColor(0x10,0x10,0x28),
         line_color=color, line_width=Pt(1))
    rect(s, Inches(6.8), y, Inches(1.5), Inches(0.4), color)
    txt(s, name, Inches(6.8), y + Inches(0.05), Inches(1.5), Inches(0.32),
        font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, line in enumerate(pros.split("\n")):
        txt(s, line, Inches(7.0), y + Inches(0.55 + j * 0.3), Inches(5.8), Inches(0.3),
            font_size=Pt(11), color=WHITE_70)

# ═══════════════════════════════════════════════════════
# SLIDE 12: 现状诊断
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "现状诊断 · 体验痛点", "大众官网 / 小程序 / App 爱车 专家走查", tag="02 背景分析")

pain_points = [
    ("🔴 C端 - 官网", [
        "信息架构层级过深，潜客找车效率低",
        "移动端响应式适配不足，主体构图错位",
        "视觉风格老旧，与 ID.ERA 科技感脱节",
        "CTA 转化路径不清晰，试驾转化率偏低",
    ], ACCENT_BLUE),
    ("🟡 C端 - 小程序", [
        "各平台小程序功能不统一，体验割裂",
        "爱车功能入口分散，用户使用频率低",
        "运营活动与常规功能界面语言不一致",
        "缺乏情感化设计，品牌感感知薄弱",
    ], ACCENT_LIGHT),
    ("🟠 C端 - App 爱车", [
        "爱车页视觉陈旧，急需全面升级",
        "车辆控制入口层级过深",
        "缺乏智能推荐与主动服务触点",
        "用户增长/留存机制设计不完善",
    ], ACCENT_GRAD),
    ("🔵 C车 - 车控需求", [
        "需全新设计CMP21D专属App",
        "大众×小鹏双品牌融合视觉有待定义",
        "遥控泊车/智驾考试等复杂交互需专业设计",
        "手机+手表双端适配体验需统一规范",
    ], TEAL),
]

for i, (title, points, color) in enumerate(pain_points):
    col = i % 2
    row = i // 2
    x = Inches(0.4 + col * 6.4)
    y = Inches(1.65 + row * 2.6)
    rect(s, x, y, Inches(6.1), Inches(2.4), RGBColor(0x10,0x10,0x28),
         line_color=color, line_width=Pt(1))
    rect(s, x, y, Inches(6.1), Inches(0.45), color)
    txt(s, title, x + Inches(0.15), y + Inches(0.06), Inches(5.8), Inches(0.35),
        font_size=Pt(14), bold=True, color=WHITE)
    for j, pt in enumerate(points):
        txt(s, "▸ " + pt, x + Inches(0.2), y + Inches(0.55 + j * 0.44), Inches(5.8), Inches(0.4),
            font_size=Pt(11.5), color=WHITE_70)

# ═══════════════════════════════════════════════════════
# SLIDE 13: 章节分隔 - 03A C端
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
rect(s, 0, 0, Inches(5.5), H, RGBColor(0x00,0x44,0xAA))
txt(s, "03A", Inches(0.5), Inches(1.8), Inches(4.5), Inches(2),
    font_size=Pt(96), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "C端产品设计方案", Inches(5.8), Inches(2.7), Inches(7.3), Inches(0.9),
    font_size=Pt(32), bold=True, color=WHITE)
txt(s, "官网 · 小程序 · App 爱车 · 视觉风格",
    Inches(5.8), Inches(3.7), Inches(7.3), Inches(0.6),
    font_size=Pt(17), color=WHITE_70)
txt(s, "C-Touch Product  |  ID.ERA Digital Experience Reinvention",
    Inches(5.8), Inches(4.4), Inches(7.3), Inches(0.5),
    font_size=Pt(13), color=WHITE_40, italic=True)
divider_line(s, Inches(5.8), Inches(2.5), Inches(7), ACCENT_BLUE, Pt(3))

# ═══════════════════════════════════════════════════════
# SLIDE 14: C端设计策略
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "C端设计策略", "以 ID.ERA 为核心，重构大众中国 C 端数字体验生态", tag="03A C端方案")

txt(s, '"以 ID.ERA 为核心，重构大众中国 C 端数字体验生态"',
    Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.55),
    font_size=Pt(19), bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

strategies = [
    ("沉浸式体验", "Immersive", "一页式滚动叙事\n场景分区构图\n视差动效吸引力", ACCENT_BLUE, "🎬"),
    ("响应式多端", "Responsive", "PC + Mobile 主体置中\n多平台统一设计语言\nDesign System 规范化", ACCENT_LIGHT, "📱"),
    ("情感化旅程", "Emotional", "三类用户差异化旅程\nMOT 关键时刻设计\n品牌温度感植入", ACCENT_GRAD, "❤️"),
    ("AI 赋能运营", "AI-Driven", "智能内容推荐\n个性化展示逻辑\n数据驱动持续优化", GOLD, "🤖"),
]

for i, (title, en, desc, color, icon) in enumerate(strategies):
    x = Inches(0.4 + i * 3.2)
    y = Inches(2.4)
    rect(s, x, y, Inches(3.0), Inches(3.8), RGBColor(0x0E,0x0E,0x28),
         line_color=color, line_width=Pt(1.5))
    rect(s, x, y, Inches(3.0), Inches(0.55), color)
    txt(s, icon + " " + title, x + Inches(0.15), y + Inches(0.08), Inches(2.7), Inches(0.4),
        font_size=Pt(15), bold=True, color=WHITE)
    txt(s, en, x + Inches(0.15), y + Inches(0.62), Inches(2.7), Inches(0.35),
        font_size=Pt(12), color=color, italic=True)
    divider_line(s, x + Inches(0.15), y + Inches(1.05), Inches(2.7), color, Pt(0.5))
    for j, line in enumerate(desc.split("\n")):
        txt(s, "▸ " + line, x + Inches(0.15), y + Inches(1.2 + j * 0.55), Inches(2.7), Inches(0.48),
            font_size=Pt(12), color=WHITE_70)

# ═══════════════════════════════════════════════════════
# SLIDE 15: 用户旅程重构
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "三类用户旅程重构", "潜在客户 · 存量客户 · 增量用户 · 全生命周期覆盖", tag="03A C端方案")

user_types = [
    ("潜在客户", "Prospect", ACCENT_BLUE, [
        ("认知 → 品牌触达", "官网沉浸首页 · 全系车型浏览"),
        ("兴趣 → 了解产品", "ID.ERA 车型详情 · 金融计算器"),
        ("决策 → 预约行动", "预约试驾 · 经销商查询 · CTA"),
    ]),
    ("存量客户", "Owner", ACCENT_GRAD, [
        ("交付 → 用车服务", "爱车App · 车控 · 充电服务"),
        ("维保 → 售后体验", "预约维保 · 爱车课堂 · 健康报告"),
        ("口碑 → 粘性留存", "官方社区 · 大众一家 · 积分体系"),
    ]),
    ("增量用户", "Growth", GOLD, [
        ("拉新 → 品牌扩散", "营销活动H5 · 分享裂变机制"),
        ("促活 → 功能深用", "好物商城 · 运营活动 · 任务系统"),
        ("转化 → 增购/置换", "官方二手车 · 品质保障 · 置换流程"),
    ]),
]

for col, (utype, en, color, stages) in enumerate(user_types):
    x = Inches(0.4 + col * 4.3)
    rect(s, x, Inches(1.65), Inches(4.05), Inches(5.5), RGBColor(0x0C,0x0C,0x22),
         line_color=color, line_width=Pt(1.5))
    rect(s, x, Inches(1.65), Inches(4.05), Inches(0.5), color)
    txt(s, utype, x + Inches(0.15), Inches(1.68), Inches(2), Inches(0.42),
        font_size=Pt(16), bold=True, color=WHITE)
    txt(s, en, x + Inches(2.2), Inches(1.7), Inches(1.7), Inches(0.38),
        font_size=Pt(11), color=WHITE_40, italic=True)
    for i, (stage, desc) in enumerate(stages):
        sy = Inches(2.3 + i * 1.5)
        rect(s, x + Inches(0.2), sy, Inches(3.6), Inches(1.2), RGBColor(0x14,0x14,0x30),
             line_color=color, line_width=Pt(0.5))
        rect(s, x + Inches(0.2), sy, Inches(3.6), Inches(0.06), color)
        txt(s, stage, x + Inches(0.35), sy + Inches(0.1), Inches(3.2), Inches(0.38),
            font_size=Pt(13), bold=True, color=color)
        txt(s, desc, x + Inches(0.35), sy + Inches(0.55), Inches(3.2), Inches(0.55),
            font_size=Pt(11), color=WHITE_70, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 16: C端信息架构 & 功能范围
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "C端产品范围 · 信息架构", "官网 · 小程序 · App 爱车 全矩阵覆盖", tag="03A C端方案")

products = [
    ("官网\n(PC + Mobile)", ACCENT_BLUE, [
        "首页 / 全系车型 / 车辆详情",
        "ID.生而电动 / ID.科技 / ID.充电",
        "官方商城 / 售后服务 / 二手车",
        "大众一家 / 经销商 / 预约试驾",
    ]),
    ("小程序\n(微信/百度/抖音等)", ACCENT_LIGHT, [
        "发现 / 服务 / 爱车 / 我的 / 好物",
        "多平台（微信/快手/小红书/百度）",
        "营销活动H5 · 刮奖/众筹/V-Day",
        "统一设计语言跨平台适配",
    ]),
    ("App 爱车\n(手机 + 手表)", ACCENT_GRAD, [
        "快捷卡片 / 车图 / 车辆控制",
        "远程空调 / 智能座舱 / 充电服务",
        "用车服务 / 官方二手车 / 车载权益",
        "智慧车联 / 爱车维修 / 全页视觉升级",
    ]),
    ("Design System", GOLD, [
        "统一品牌图腾 / Monogram",
        "色彩/字体/间距规范",
        "组件库 · Icon · 动效规范",
        "多端复用 · 持续迭代指导",
    ]),
]

for i, (name, color, items) in enumerate(products):
    x = Inches(0.4 + i * 3.2)
    rect(s, x, Inches(1.65), Inches(3.0), Inches(5.5), RGBColor(0x0E,0x0E,0x26),
         line_color=color, line_width=Pt(1))
    rect(s, x, Inches(1.65), Inches(3.0), Inches(0.55), color)
    txt(s, name, x + Inches(0.15), Inches(1.68), Inches(2.7), Inches(0.5),
        font_size=Pt(13), bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, "▸ " + item, x + Inches(0.15), Inches(2.35 + j * 1.1), Inches(2.7), Inches(0.95),
            font_size=Pt(11.5), color=WHITE_70, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 17: C端视觉风格方向
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "C端视觉风格方向", "Look & Feel · ID.ERA 品牌数字调性探索", tag="03A C端方案")

txt(s, "方向 A：科技电动 · ID.ERA 蓝",
    Inches(0.5), Inches(1.65), Inches(4), Inches(0.45),
    font_size=Pt(14), bold=True, color=ACCENT_BLUE)
rect(s, Inches(0.5), Inches(2.15), Inches(3.8), Inches(4.5), RGBColor(0x04,0x0A,0x1E),
     line_color=ACCENT_BLUE, line_width=Pt(1.5))
for i, (label, color) in enumerate([
    ("主色 Science Blue #008BFF", ACCENT_BLUE),
    ("辅色 Electric Cyan #00D4AA", ACCENT_GRAD),
    ("背景 Deep Space #0A0A14", BG_DARK),
    ("高光 Pure White #FFFFFF", WHITE),
]):
    rect(s, Inches(0.7), Inches(2.35 + i * 0.52), Inches(0.4), Inches(0.35), color)
    txt(s, label, Inches(1.2), Inches(2.38 + i * 0.52), Inches(3.0), Inches(0.32),
        font_size=Pt(11), color=WHITE_70)
txt(s, "特点：科技感强 · 与ID.ERA电动属性高度契合\n适用：官网/App主色调方向",
    Inches(0.7), Inches(4.5), Inches(3.4), Inches(0.6),
    font_size=Pt(11), color=WHITE_40, wrap=True)

txt(s, "方向 B：暖科技 · 中国温度",
    Inches(4.7), Inches(1.65), Inches(4), Inches(0.45),
    font_size=Pt(14), bold=True, color=GOLD)
rect(s, Inches(4.7), Inches(2.15), Inches(3.8), Inches(4.5), RGBColor(0x14,0x0C,0x04),
     line_color=GOLD, line_width=Pt(1.5))
for i, (label, color) in enumerate([
    ("主色 Warm Gold #FFC800", GOLD),
    ("辅色 Sunset Orange #FF6B35", RGBColor(0xFF,0x6B,0x35)),
    ("背景 Warm Night #14100A", RGBColor(0x14,0x10,0x0A)),
    ("强调 Pearl White #F8F4EF", RGBColor(0xF8,0xF4,0xEF)),
]):
    rect(s, Inches(4.9), Inches(2.35 + i * 0.52), Inches(0.4), Inches(0.35), color)
    txt(s, label, Inches(5.4), Inches(2.38 + i * 0.52), Inches(3.0), Inches(0.32),
        font_size=Pt(11), color=WHITE_70)
txt(s, "特点：品牌温度感强 · 中国文化融合\n适用：官网营销活动/社区方向",
    Inches(4.9), Inches(4.5), Inches(3.4), Inches(0.6),
    font_size=Pt(11), color=WHITE_40, wrap=True)

txt(s, "方向 C：极简沉浸 · 大场域构图",
    Inches(8.9), Inches(1.65), Inches(4), Inches(0.45),
    font_size=Pt(14), bold=True, color=WHITE)
rect(s, Inches(8.9), Inches(2.15), Inches(4.0), Inches(4.5), RGBColor(0x08,0x08,0x08),
     line_color=WHITE_40, line_width=Pt(1.5))
for i, (label, color) in enumerate([
    ("主色 Pure Black #080808", RGBColor(0x08,0x08,0x08)),
    ("强调 Ice White #F0F4FF", RGBColor(0xF0,0xF4,0xFF)),
    ("点缀 Chrome Blue #007AFF", RGBColor(0x00,0x7A,0xFF)),
    ("材质 Glass Blur 毛玻璃", RGBColor(0x30,0x30,0x50)),
]):
    rect(s, Inches(9.1), Inches(2.35 + i * 0.52), Inches(0.4), Inches(0.35), color,
         line_color=WHITE_40, line_width=Pt(0.5))
    txt(s, label, Inches(9.6), Inches(2.38 + i * 0.52), Inches(3.1), Inches(0.32),
        font_size=Pt(11), color=WHITE_70)
txt(s, "特点：Apple风极简·大场域视差·高级感\n适用：官网车型详情页方向",
    Inches(9.1), Inches(4.5), Inches(3.6), Inches(0.6),
    font_size=Pt(11), color=WHITE_40, wrap=True)

# 推荐标记
tag_badge(s, "★ 推荐", Inches(0.5), Inches(6.5), ACCENT_BLUE, WHITE, Pt(12))

# ═══════════════════════════════════════════════════════
# SLIDE 18: 章节分隔 - 03B C车
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, TEAL)
rect(s, 0, 0, Inches(5.5), H, RGBColor(0x00,0x3D,0x2E))
txt(s, "03B", Inches(0.5), Inches(1.8), Inches(4.5), Inches(2),
    font_size=Pt(96), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "C车设计方案", Inches(5.8), Inches(2.7), Inches(7.3), Inches(0.9),
    font_size=Pt(32), bold=True, color=WHITE)
txt(s, "CMP21D 车型移动互联 App · 全场景情感化车控",
    Inches(5.8), Inches(3.7), Inches(7.3), Inches(0.6),
    font_size=Pt(17), color=WHITE_70)
txt(s, "VW × XPENG  |  Smart Vehicle Control Experience Design",
    Inches(5.8), Inches(4.4), Inches(7.3), Inches(0.5),
    font_size=Pt(13), color=WHITE_40, italic=True)
divider_line(s, Inches(5.8), Inches(2.5), Inches(7), TEAL, Pt(3))

# ═══════════════════════════════════════════════════════
# SLIDE 19: C车设计策略
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, TEAL)
section_header(s, "C车设计策略", "全场景情感化车控 · 车随心动", tag="03B C车方案")

txt(s, '"全场景情感化车控体验 · 车随心动"',
    Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.55),
    font_size=Pt(19), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# 三大原则
principles = [
    ("01 简洁优先", "Simple First",
     "核心功能一目了然，减少点击层级\n降低学习成本，提升操作效率", ACCENT_BLUE),
    ("02 反馈及时", "Instant Feedback",
     "操作结果清晰传达，避免用户焦虑\n实时状态同步，提升系统信任感", TEAL),
    ("03 操作流畅", "Smooth Flow",
     "步骤精简连贯，减少操作卡顿\n舒适便捷体验，情感化触点设计", GOLD),
]

for i, (title, en, desc, color) in enumerate(principles):
    x = Inches(0.4 + i * 4.3)
    rect(s, x, Inches(2.35), Inches(4.0), Inches(1.6), RGBColor(0x0C,0x0C,0x22),
         line_color=color, line_width=Pt(1.5))
    rect(s, x, Inches(2.35), Inches(4.0), Inches(0.5), color)
    txt(s, title, x + Inches(0.15), Inches(2.38), Inches(3.7), Inches(0.42),
        font_size=Pt(15), bold=True, color=WHITE)
    txt(s, en, x + Inches(0.15), Inches(2.92), Inches(3.7), Inches(0.3),
        font_size=Pt(11), color=color, italic=True)
    for j, line in enumerate(desc.split("\n")):
        txt(s, line, x + Inches(0.15), Inches(3.3 + j * 0.3), Inches(3.7), Inches(0.28),
            font_size=Pt(11.5), color=WHITE_70)

# 双品牌融合策略
rect(s, Inches(0.4), Inches(4.2), Inches(12.5), Inches(2.85), RGBColor(0x08,0x18,0x18))
txt(s, "大众 × 小鹏 双品牌融合设计策略",
    Inches(0.6), Inches(4.35), Inches(12), Inches(0.45),
    font_size=Pt(16), bold=True, color=TEAL)
divider_line(s, Inches(0.6), Inches(4.85), Inches(12), TEAL, Pt(0.5))

fusion_items = [
    ("大众基因", "可信赖 · 精准 · 品质感 · 工程理性", ACCENT_BLUE),
    ("小鹏基因", "智能化 · 科技感 · AI-Native · 年轻化", RGBColor(0x00,0xCC,0x88)),
    ("融合设计语言", "沉稳底色 + 科技亮色 · 严谨结构 + 流动动效", GOLD),
]
for i, (title, desc, color) in enumerate(fusion_items):
    x = Inches(0.6 + i * 4.2)
    txt(s, title, x, Inches(5.0), Inches(4.0), Inches(0.38),
        font_size=Pt(14), bold=True, color=color)
    txt(s, desc, x, Inches(5.42), Inches(4.0), Inches(0.55),
        font_size=Pt(11.5), color=WHITE_70, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 20: 车控首页设计
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, TEAL)
section_header(s, "车控首页 Home Dashboard", "3D车模 · 实时状态 · 快捷控制 · Avatar 入口", tag="03B C车方案")

# 左侧 UI 线框示意
rect(s, Inches(0.5), Inches(1.65), Inches(5.5), Inches(5.5), RGBColor(0x0C,0x0C,0x22),
     line_color=TEAL, line_width=Pt(1.5))

# 手机框模拟
rect(s, Inches(1.5), Inches(1.8), Inches(3.5), Inches(5.2), RGBColor(0x08,0x08,0x18),
     line_color=RGBColor(0x30,0x30,0x50), line_width=Pt(1.5))

# 状态栏区域
rect(s, Inches(1.5), Inches(1.8), Inches(3.5), Inches(0.4), RGBColor(0x0A,0x0A,0x1E))
txt(s, "9:41  CMP21D", Inches(1.6), Inches(1.83), Inches(3.2), Inches(0.32),
    font_size=Pt(9), color=WHITE_40)

# 车辆信息区
rect(s, Inches(1.5), Inches(2.2), Inches(3.5), Inches(1.5), RGBColor(0x06,0x12,0x22))
txt(s, "CMP21D  |  沪 A 88888", Inches(1.65), Inches(2.25), Inches(3.2), Inches(0.35),
    font_size=Pt(11), bold=True, color=WHITE)
txt(s, "⚡ 82%  ·  续航 389km  ·  已锁车", Inches(1.65), Inches(2.6), Inches(3.2), Inches(0.35),
    font_size=Pt(10), color=TEAL)

# 3D车模占位
rect(s, Inches(1.65), Inches(2.98), Inches(3.2), Inches(1.2), RGBColor(0x08,0x16,0x2A),
     line_color=ACCENT_BLUE, line_width=Pt(0.5))
txt(s, "[ 3D 车模  可交互旋转 ]", Inches(1.65), Inches(3.3), Inches(3.2), Inches(0.4),
    font_size=Pt(11), color=WHITE_40, align=PP_ALIGN.CENTER)

# 快捷控制区
quick_btns = ["🔒解锁", "❄️空调", "📍定位", "💡闪灯"]
for j, btn in enumerate(quick_btns):
    bx = Inches(1.65 + j * 0.82)
    rect(s, bx, Inches(4.28), Inches(0.72), Inches(0.55), RGBColor(0x14,0x22,0x3A),
         line_color=ACCENT_BLUE, line_width=Pt(0.5))
    txt(s, btn, bx, Inches(4.32), Inches(0.72), Inches(0.42),
        font_size=Pt(9), color=WHITE_70, align=PP_ALIGN.CENTER)

# OneHit 场景区
rect(s, Inches(1.65), Inches(4.95), Inches(3.2), Inches(0.45), RGBColor(0x10,0x20,0x38))
txt(s, "🌟 OneHit  上班模式  · 旅行模式  · 放松模式",
    Inches(1.68), Inches(4.98), Inches(3.1), Inches(0.38),
    font_size=Pt(9), color=ACCENT_LIGHT)

# Avatar 入口
rect(s, Inches(1.65), Inches(5.5), Inches(3.2), Inches(0.38), RGBColor(0x00,0x22,0x44))
txt(s, "🤖 Avatar  晚上好，上海今晚 8°C，已预热座椅",
    Inches(1.68), Inches(5.53), Inches(3.1), Inches(0.32),
    font_size=Pt(9), color=ACCENT_GRAD)

# 右侧 UX 说明
ux_items = [
    ("车辆状态概览", "电量/续航/车身状态/温度实时显示", TEAL),
    ("快捷控制区", "远程解闭锁 · 空调开关 · 闪灯鸣笛", ACCENT_BLUE),
    ("OneHit 场景", "智能推荐模式 · 一键激活联动设置", ACCENT_LIGHT),
    ("Avatar 入口", "情感化欢迎语 · 大模型智慧服务", ACCENT_GRAD),
    ("OTA 通知栏", "升级进度 · 哨兵模式警报推送", GOLD),
]
for i, (feature, desc, color) in enumerate(ux_items):
    y = Inches(1.75 + i * 1.02)
    rect(s, Inches(6.3), y, Inches(6.6), Inches(0.85), RGBColor(0x0E,0x0E,0x26),
         line_color=color, line_width=Pt(0.8))
    rect(s, Inches(6.3), y, Inches(0.06), Inches(0.85), color)
    txt(s, feature, Inches(6.5), y + Inches(0.08), Inches(6.2), Inches(0.35),
        font_size=Pt(13), bold=True, color=color)
    txt(s, desc, Inches(6.5), y + Inches(0.47), Inches(6.2), Inches(0.32),
        font_size=Pt(11), color=WHITE_40)

# ═══════════════════════════════════════════════════════
# SLIDE 21: 核心功能 UX 设计
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, TEAL)
section_header(s, "核心功能 UX 设计", "车控 · 空调 · 遥控泊车 · 充电 · 智驾", tag="03B C车方案")

features = [
    ("远程空调\n& 一键备车", [
        "圆弧滑动条 · 温度实时跳动",
        "极速升/降温 · 风速风向控制",
        "座椅加热/通风联动",
        "状态反馈：正在降温至24℃",
    ], ACCENT_BLUE),
    ("遥控泊车\n& 辅助驾驶", [
        "车位扫描实时显示",
        "遥控泊入/泊出操控",
        "周边环境 + 障碍物可视化",
        "安全提醒全程实时反馈",
    ], ACCENT_GRAD),
    ("充电服务\n& 电池管理", [
        "充电启停控制 · 截止电量设置",
        "预约充电 · 交流电流限值",
        "电池预热 · 远程放电控制",
        "充电进度可视化动效",
    ], TEAL),
    ("数字钥匙\n& 权限管理", [
        "分享钥匙 · 权限分级 P/S 用户",
        "时间段授权 · 远程取消授权",
        "授权记录追踪",
        "安全提醒 · 操作二次确认",
    ], GOLD),
    ("智驾考试\n& 功能解锁", [
        "扫码跳转视频学习页",
        "在线考试 · 通过即解锁智驾",
        "绑定P/S用户权限体系",
        "学习进度可视化追踪",
    ], ACCENT_LIGHT),
    ("OneHit 情景模式", [
        "场景乐高编辑器 · 拖拽组合",
        "温度/香氛/灯光/导航联动",
        "场景广场 · 公开共享",
        "智能预测：基于行程推荐模式",
    ], RED_LIGHT),
]

for i, (title, items, color) in enumerate(features):
    col = i % 3
    row = i // 2
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.65 + row * 2.7)
    rect(s, x, y, Inches(4.0), Inches(2.45), RGBColor(0x0C,0x0C,0x22),
         line_color=color, line_width=Pt(1))
    rect(s, x, y, Inches(4.0), Inches(0.5), color)
    txt(s, title, x + Inches(0.15), y + Inches(0.06), Inches(3.7), Inches(0.42),
        font_size=Pt(13), bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, "▸ " + item, x + Inches(0.18), y + Inches(0.62 + j * 0.43), Inches(3.65), Inches(0.4),
            font_size=Pt(11), color=WHITE_70)

# ═══════════════════════════════════════════════════════
# SLIDE 22: C车视觉风格
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, TEAL)
section_header(s, "C车视觉风格方向", "大众 × 小鹏 双品牌融合视觉探索", tag="03B C车方案")

# 主推方向
txt(s, "主推方向：「深空智控」Dark Intelligence",
    Inches(0.5), Inches(1.65), Inches(12), Inches(0.45),
    font_size=Pt(16), bold=True, color=TEAL)

# 色板
palette = [
    ("Deep Navy #0A0A14", BG_DARK, "主背景"),
    ("VW Blue #0033CC", RGBColor(0x00,0x33,0xCC), "大众品牌蓝"),
    ("XP Teal #00C9A7", TEAL, "小鹏科技青"),
    ("Alert Gold #FFC800", GOLD, "操作强调金"),
    ("Pure White #FFFFFF", WHITE, "文字/图标"),
    ("Warning Red #FF4D6D", RED_LIGHT, "警告状态"),
]
for i, (name, color, usage) in enumerate(palette):
    x = Inches(0.5 + i * 2.13)
    rect(s, x, Inches(2.2), Inches(1.95), Inches(1.1), color,
         line_color=RGBColor(0x30,0x30,0x50), line_width=Pt(0.5))
    txt(s, usage, x, Inches(3.38), Inches(1.95), Inches(0.28),
        font_size=Pt(10), color=WHITE_40, align=PP_ALIGN.CENTER)
    txt(s, name.split(" ")[0], x, Inches(3.68), Inches(1.95), Inches(0.28),
        font_size=Pt(9), color=WHITE_40, align=PP_ALIGN.CENTER)

# 设计元素说明
design_els = [
    ("车控页面", "3D车模居中 · 深色背景 · 状态信息叠加 · 毛玻璃卡片", ACCENT_BLUE),
    ("交互动效", "按钮点击波纹 · 解锁车门高亮 · 空调温度波动 · OTA进度条", TEAL),
    ("图标体系", "扁平化设计 · 线面结合 · 车控专属图标集 · 统一视觉重量", ACCENT_GRAD),
    ("手表端适配", "极简信息架构 · 表冠交互 · 单手操作优先 · 圆形构图", GOLD),
]
for i, (title, desc, color) in enumerate(design_els):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(4.35 + row * 1.4)
    rect(s, x, y, Inches(6.1), Inches(1.1), RGBColor(0x0E,0x0E,0x26),
         line_color=color, line_width=Pt(0.8))
    rect(s, x, y, Inches(0.06), Inches(1.1), color)
    txt(s, title, x + Inches(0.2), y + Inches(0.1), Inches(5.7), Inches(0.35),
        font_size=Pt(13), bold=True, color=color)
    txt(s, desc, x + Inches(0.2), y + Inches(0.55), Inches(5.7), Inches(0.45),
        font_size=Pt(11), color=WHITE_40, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 23: 章节分隔 - 04 项目计划
# ═══════════════════════════════════════════════════════
s = add_slide()
chapter_divider(s, 4, "项目计划", "Team & Plan & Deliverables  ·  排期 & 团队 & 交付物")

# ═══════════════════════════════════════════════════════
# SLIDE 24: 项目团队
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "项目实施团队", "全岗位配置 · 对应 RFP 核心人员要求", tag="04 项目计划")

# 核心三人
core = [
    ("Elya", "设计总监 / 项目经理", "U9 · 14年体验设计\n前百度云UED负责人\n大众·奥迪·保时捷项目负责人", ACCENT_BLUE),
    ("Maya", "创意总监 / Design Leader", "U8 · 15年产品设计\n前三星·百度MUX资深设计\n保时捷·林肯·现代项目负责人", ACCENT_LIGHT),
    ("Arthur", "策略总监 / Strategy Leader", "P9 · 15年顶层设计\n前独角兽用户增长负责人\n极氪·吉利·奥迪全链路专家", ACCENT_GRAD),
]
for i, (name, role, exp, color) in enumerate(core):
    x = Inches(0.4 + i * 4.3)
    rect(s, x, Inches(1.65), Inches(4.05), Inches(2.0), RGBColor(0x0A,0x0A,0x22),
         line_color=color, line_width=Pt(2))
    rect(s, x, Inches(1.65), Inches(4.05), Inches(0.5), color)
    txt(s, name, x + Inches(0.2), Inches(1.68), Inches(1.5), Inches(0.42),
        font_size=Pt(18), bold=True, color=WHITE)
    txt(s, role, x + Inches(1.7), Inches(1.72), Inches(2.2), Inches(0.38),
        font_size=Pt(11), color=color)
    for j, line in enumerate(exp.split("\n")):
        txt(s, "· " + line, x + Inches(0.2), Inches(2.28 + j * 0.32), Inches(3.7), Inches(0.3),
            font_size=Pt(11), color=WHITE_40)

# 执行团队
txt(s, "执行设计团队", Inches(0.4), Inches(3.85), Inches(3), Inches(0.4),
    font_size=Pt(14), bold=True, color=WHITE)
exec_team = [
    ("James", "UI 设计总监", "2021红点奖获得者", TEAL),
    ("Kay", "高级交互设计师", "11年从业·车联网专家", ACCENT_BLUE),
    ("Elisa", "高级视觉设计师", "车厂App视觉专家", ACCENT_LIGHT),
    ("PIXELSIR", "动效设计师", "场景动效·微交互专家", ACCENT_GRAD),
    ("Mira", "用户研究员", "可用性测试·用户访谈", GOLD),
    ("Dandan", "用户体验专家", "数据分析·体验评估", WHITE_70),
]
for i, (name, role, spec, color) in enumerate(exec_team):
    x = Inches(0.4 + i * 2.15)
    rect(s, x, Inches(4.35), Inches(2.0), Inches(2.5), RGBColor(0x0E,0x0E,0x26),
         line_color=color, line_width=Pt(0.8))
    rect(s, x, Inches(4.35), Inches(2.0), Inches(0.06), color)
    txt(s, name, x + Inches(0.1), Inches(4.45), Inches(1.8), Inches(0.38),
        font_size=Pt(13), bold=True, color=color)
    txt(s, role, x + Inches(0.1), Inches(4.85), Inches(1.8), Inches(0.35),
        font_size=Pt(11), bold=True, color=WHITE)
    txt(s, spec, x + Inches(0.1), Inches(5.25), Inches(1.8), Inches(0.55),
        font_size=Pt(10), color=WHITE_40, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 25: 项目时间表
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "项目时间表", "C端 & C车 双项目并行推进", tag="04 项目计划")

# 时间轴
timeline = [
    ("11.07", "RFP 发出", ACCENT_BLUE, True),
    ("11.14", "提案提交", ACCENT_LIGHT, True),
    ("11.17", "讲标确认", ACCENT_GRAD, True),
    ("11.18+", "项目启动", TEAL, False),
    ("12月", "用研·场景分析", GOLD, False),
    ("1月", "信息架构·交互", ACCENT_BLUE, False),
    ("2-3月", "视觉设计·原型", ACCENT_LIGHT, False),
    ("4-5月", "开发支持·测试", TEAL, False),
    ("6月", "上线·维护", GOLD, False),
]

# 时间线横轴
rect(s, Inches(0.6), Inches(2.6), Inches(12.3), Inches(0.04), WHITE_40)

for i, (time, label, color, done) in enumerate(timeline):
    x = Inches(0.6 + i * 1.4)
    # 节点圆点
    circle = slide.shapes if False else None
    dot = s.shapes.add_shape(9, x - Inches(0.12), Inches(2.45), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color if done else RGBColor(0x30,0x30,0x50)
    dot.line.color.rgb = color

    txt(s, time, x - Inches(0.4), Inches(1.85), Inches(0.9), Inches(0.35),
        font_size=Pt(11), bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, x - Inches(0.5), Inches(2.9), Inches(1.1), Inches(0.55),
        font_size=Pt(10), color=WHITE_70 if done else WHITE_40, align=PP_ALIGN.CENTER, wrap=True)

# 两个项目甘特
phases_cend = [
    ("C端 项目", [
        ("用户研究 & 场景分析", ACCENT_BLUE, 0, 1.5),
        ("信息架构 & 交互设计", ACCENT_LIGHT, 1.5, 3),
        ("视觉设计 & 原型", ACCENT_GRAD, 3, 5.5),
        ("开发支持 & 可用性测试", TEAL, 5.5, 7.5),
        ("持续优化 & 维护", GOLD, 7.5, 9),
    ]),
    ("C车 项目", [
        ("用户研究 & 场景分析", ACCENT_BLUE, 0, 1.5),
        ("信息架构 & 交互设计", ACCENT_LIGHT, 1.5, 3.5),
        ("视觉设计 & 原型", ACCENT_GRAD, 3.5, 5.5),
        ("开发支持 & 可用性测试", TEAL, 5.5, 7.5),
        ("持续优化 & 维护", GOLD, 7.5, 9),
    ]),
]

for row_i, (proj_name, phases) in enumerate(phases_cend):
    y = Inches(3.7 + row_i * 1.5)
    rect(s, Inches(0.6), y, Inches(1.6), Inches(1.1), RGBColor(0x10,0x10,0x28))
    txt(s, proj_name, Inches(0.65), y + Inches(0.35), Inches(1.5), Inches(0.4),
        font_size=Pt(11), bold=True, color=WHITE)
    scale = Inches(1.3)
    for (phase_name, color, start, end) in phases:
        bx = Inches(2.3) + int(start * scale)
        bw = int((end - start) * scale)
        rect(s, bx, y + Inches(0.15), bw, Inches(0.75), color)
        if bw > Inches(0.8):
            txt(s, phase_name, bx + Inches(0.08), y + Inches(0.28), bw - Inches(0.1), Inches(0.5),
                font_size=Pt(9), color=WHITE, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 26: 交付物清单
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "项目交付物清单", "完整覆盖 RFP 全阶段交付要求", tag="04 项目计划")

deliverables = [
    ("🔍 用研阶段", ACCENT_BLUE, [
        "用户画像 / Personas",
        "竞品分析报告 / Competitive Analysis",
        "服务蓝图 / Service Blueprint",
        "故事板 / Storyboard",
        "视觉风格定义 / Look & Feel",
    ]),
    ("🗺️ 场景&架构", ACCENT_LIGHT, [
        "用户旅程地图 / User Journey Map",
        "用例 / Use Cases",
        "用户流程图 / User Flow Chart",
        "信息架构图 / IA Diagram",
        "线框图 / Wireframe",
    ]),
    ("🎨 设计阶段", ACCENT_GRAD, [
        "页面交互文档 / UX Documentation",
        "UX 交互规范 / UX Guideline",
        "视觉原型 / UI Prototype",
        "切图标注 / Slicing & Marking",
        "动效文件 / Animation Demo",
    ]),
    ("🚀 落地阶段", TEAL, [
        "设计评审会纪要 / Meeting Minutes",
        "可用性测试报告 / UT Report",
        "项目成果报告 / Project Results",
        "优化方案报告 / Optimization Plan",
        "上线后1个月迭代支持",
    ]),
]

for i, (title, color, items) in enumerate(deliverables):
    x = Inches(0.4 + i * 3.2)
    rect(s, x, Inches(1.65), Inches(3.0), Inches(5.5), RGBColor(0x0C,0x0C,0x22),
         line_color=color, line_width=Pt(1))
    rect(s, x, Inches(1.65), Inches(3.0), Inches(0.5), color)
    txt(s, title, x + Inches(0.15), Inches(1.68), Inches(2.7), Inches(0.42),
        font_size=Pt(13), bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(s, "✓ " + item, x + Inches(0.15), Inches(2.3 + j * 0.96), Inches(2.7), Inches(0.85),
            font_size=Pt(11.5), color=WHITE_70, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 27: 竞争优势 & 承诺
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
section_header(s, "我们的竞争优势", "为何选择 TOPX · 我们的承诺", tag="04 项目计划")

advantages = [
    ("专业性", "深度行业经验", [
        "大众·奥迪·保时捷·斯柯达 5年+合作",
        "曾负责上汽大众C端产品设计 (加分项)",
        "车联网App UIUX设计丰富经验",
    ], ACCENT_BLUE, "🏆"),
    ("创新性", "独特设计方案", [
        "OneHit场景模式/Avatar大模型首创",
        "双品牌融合设计语言方法论",
        "5维体验诊断独家方法论体系",
    ], ACCENT_LIGHT, "💡"),
    ("执行力", "高效可靠交付", [
        "严格遵循时间节点·参与日常例会",
        "安亭现场驻厂·无缝沟通协作",
        "STAR规划·层层递进确保质量",
    ], ACCENT_GRAD, "⚡"),
    ("本土化", "深度中国理解", [
        "中国用户使用习惯深度洞察",
        "多平台生态（微信/抖音/百度）专家",
        "In China for China 设计策略",
    ], TEAL, "🇨🇳"),
]

for i, (title, sub, items, color, icon) in enumerate(advantages):
    x = Inches(0.4 + i * 3.2)
    rect(s, x, Inches(1.65), Inches(3.0), Inches(4.0), RGBColor(0x0C,0x0C,0x22),
         line_color=color, line_width=Pt(1.5))
    rect(s, x, Inches(1.65), Inches(3.0), Inches(0.55), color)
    txt(s, icon + " " + title, x + Inches(0.15), Inches(1.68), Inches(2.7), Inches(0.42),
        font_size=Pt(15), bold=True, color=WHITE)
    txt(s, sub, x + Inches(0.15), Inches(2.28), Inches(2.7), Inches(0.32),
        font_size=Pt(11), color=color, italic=True)
    for j, item in enumerate(items):
        txt(s, "▸ " + item, x + Inches(0.15), Inches(2.68 + j * 0.55), Inches(2.7), Inches(0.5),
            font_size=Pt(11), color=WHITE_70, wrap=True)

# 愿景承诺
rect(s, Inches(0.4), Inches(5.85), Inches(12.5), Inches(1.2), RGBColor(0x00,0x22,0x55))
rect(s, Inches(0.4), Inches(5.85), Inches(0.08), Inches(1.2), ACCENT_BLUE)
txt(s, "我们的承诺",
    Inches(0.65), Inches(5.9), Inches(2.5), Inches(0.38),
    font_size=Pt(14), bold=True, color=ACCENT_BLUE)
txt(s, ("「期待通过本次合作，为上汽大众 ID.ERA 与 CMP21D 打造行业领先的 C端数字体验与车控App，"
        "提升用户忠诚度与品牌价值，实现项目与品牌的双赢。」"),
    Inches(0.65), Inches(6.35), Inches(12.2), Inches(0.55),
    font_size=Pt(13), color=WHITE_70, italic=True, wrap=True)

# ═══════════════════════════════════════════════════════
# SLIDE 28: 结尾页
# ═══════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
rect(s, 0, 0, Inches(0.12), H, ACCENT_BLUE)
rect(s, 0, H - Inches(0.5), W, Inches(0.5), RGBColor(0x00,0x33,0x88))

# 顶部细线
rect(s, Inches(0.3), Inches(0.4), Inches(12.7), Inches(0.02), ACCENT_BLUE)
rect(s, Inches(0.3), H - Inches(0.55), Inches(12.7), Inches(0.02), ACCENT_BLUE)

txt(s, "谢谢", Inches(0.5), Inches(1.8), Inches(12.3), Inches(2.2),
    font_size=Pt(100), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Thanks for Watching",
    Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.6),
    font_size=Pt(22), color=ACCENT_BLUE, align=PP_ALIGN.CENTER, italic=True)

divider_line(s, Inches(4), Inches(4.55), Inches(5.3), ACCENT_BLUE, Pt(1.5))

txt(s, "TOPX Design  |  www.topx.design",
    Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.4),
    font_size=Pt(14), color=WHITE_70, align=PP_ALIGN.CENTER)
txt(s, "2025 上汽大众 C端产品 & C车 UIUX 设计提案",
    Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.4),
    font_size=Pt(13), color=WHITE_40, align=PP_ALIGN.CENTER)
txt(s, "CONFIDENTIAL  |  For SVW CIX Only",
    Inches(0.5), H - Inches(0.45), Inches(12.3), Inches(0.35),
    font_size=Pt(10), color=WHITE_40, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# 保存文件
# ─────────────────────────────────────────────
output_path = '/Users/ccelya/Downloads/AI-SVW-C/Proposal/SVW-TOPX-提案-2025.pptx'
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")
